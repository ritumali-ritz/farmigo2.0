from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# --- Configurations ---
OUTPUT_FILE = "Farmiga_Project_Presentation.pptx"
LOGO_PATH = "web/Green Orange Illustration Farm Logo.png" # Adjust if needed
TITLE = "Farmiga: Farm to Table App"
SUBTITLE = "Bridging the Gap Between Farmers and Consumers"
AUTHOR = "Developed by Ritex Studios"

# Theme Colors
PRIMARY_COLOR = RGBColor(46, 125, 50) # #2E7D32 (Deep Green)
ACCENT_COLOR = RGBColor(129, 199, 132) # #81C784 (Light Green)
TEXT_COLOR = RGBColor(33, 33, 33) # #212121

# --- Helper Functions ---
def add_slide(prs, layout_index, title_text, content_text=None):
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set Title
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        title_format = slide.shapes.title.text_frame.paragraphs[0].font
        title_format.name = 'Arial'
        title_format.size = Pt(36)
        title_format.bold = True
        title_format.color.rgb = PRIMARY_COLOR

    # Set Content
    if content_text and len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        body.text = content_text
        for paragraph in body.text_frame.paragraphs:
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = TEXT_COLOR
            
    return slide

def add_bullet_slide(prs, title, bullets):
    slide = add_slide(prs, 1, title) # 1 = Title and Content
    body = slide.placeholders[1]
    body.text = "" # Clear default
    tf = body.text_frame
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(22)
        p.level = 0
        p.font.name = 'Arial'
        
    return slide

def add_section_header(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[2]) # 2 = Section Header
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    return slide

def add_image_slide(prs, title, image_text_placeholder):
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # 5 = Title Only
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Add a placeholder box for the image
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4.5)
    
    shape = slide.shapes.add_shape(
        1, left, top, width, height # 1 = Rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
    shape.line.color.rgb = ACCENT_COLOR
    
    text_frame = shape.text_frame
    text_frame.text = image_text_placeholder
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return slide

# --- Main Generation ---
prs = Presentation()

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = TITLE
subtitle.text = f"{SUBTITLE}\n\n{AUTHOR}"

title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
title.text_frame.paragraphs[0].font.bold = True

# Slide 2: Introduction
add_bullet_slide(prs, "Introduction", [
    "Farmiga is a direct-to-consumer marketplace app.",
    "It connects local farmers directly with buyers.",
    "Eliminates middlemen, ensuring fresher produce and better prices.",
    "Built with modern cross-platform technology."
])

# Slide 3: Problem Statement
add_bullet_slide(prs, "The Problem", [
    "Farmers get low prices due to middlemen.",
    "Consumers pay high prices for stale produce.",
    "Lack of transparency in the supply chain.",
    "Difficulty in finding organic, local food sources."
])

# Slide 4: The Solution
add_bullet_slide(prs, "The Solution: Farmiga", [
    "A unified platform for Farmers and Buyers.",
    "Farmers list products with their own prices.",
    "Buyers order fresh produce for home delivery.",
    "Transparent pricing and direct connection.",
    "Rating and review system for trust."
])

# Slide 5: Tech Stack
add_bullet_slide(prs, "Technology Stack", [
    "Frontend: Flutter (Dart) - Cross-platform (Android/iOS/Web).",
    "Backend: Firebase (Auth, Firestore DB).",
    "Notifications: Flutter Local Notifications & FCM.",
    "Architecture: MVVM with Provider State Management.",
    "IDE: VS Code."
])

# Slide 6: Key Features
add_bullet_slide(prs, "Key Features", [
    "Role-based Authentication (Farmer/Buyer).",
    "Smart Product Image Matching.",
    "Real-time Order Notifications.",
    "Dynamic Promotional Banners.",
    "Order Tracking (Active/Past history).",
    "Geolocation for Address."
])

# Slide 7: User Roles
add_bullet_slide(prs, "User Roles", [
    "Buyer:",
    "  - Browse categories (Fruits, Veg, Dairy).",
    "  - Add to Cart & Checkout.",
    "  - Track Order Status.",
    "Farmer:",
    "  - Manage Inventory (Add/Edit/Delete products).",
    "  - View & Manage Incoming Orders.",
    "  - Business Insights."
])

# Slide 8: UI Showcase - Buyer Home
add_image_slide(prs, "Buyer Home Screen", "[INSERT SCREENSHOT: Buyer Home with Banners & Categories]")
# Note: User should replace this with actual screenshot

# Slide 9: Feature Spotlight - Banners
add_bullet_slide(prs, "Feature: Dynamic Banners", [
    "Carousel Slider with clean animations.",
    "Custom artistic gradients & shapes.",
    "Support for real promotional images.",
    "Auto-play with page indicators.",
    "Showcases offers and seasonal products."
])

# Slide 10: Feature Spotlight - Smart Matching
add_bullet_slide(prs, "Feature: Smart Image Matching", [
    "Solves the 'product photo' problem for farmers.",
    "Automatically suggests high-quality images based on keywords.",
    "Supports Hindi & English terms (e.g., 'Aloo', 'Potato').",
    "Uses curated local asset library for speed.",
    "Ensures a consistent, professional catalog look."
])

# Slide 11: Application Flow - Buying
add_image_slide(prs, "Buying Process", "[INSERT SCREENSHOT: Product Details -> Cart -> Order Success]")

# Slide 12: Application Flow - Selling
add_image_slide(prs, "Selling Process", "[INSERT SCREENSHOT: Add Product Screen -> Farmer Dashboard]")

# Slide 13: Real-time Notifications
add_bullet_slide(prs, "Tech Deep Dive: Notifications", [
    "Instant feedback for critical actions.",
    "Buyer: 'Order Placed Successfully!'",
    "Farmer: 'New Order Received!' (even in background).",
    "Implemented using flutter_local_notifications.",
    "No expensive backend server required."
])

# Slide 14: Order Management
add_bullet_slide(prs, "Order Management System", [
    "Split view for 'Active' and 'Past' orders.",
    "Status tracking: Pending -> Accepted -> Shipped -> Delivered.",
    "Client-side sorting ensures latest orders appear first.",
    "Prevents order loss and confusion."
])

# Slide 15: Profile & Settings
add_bullet_slide(prs, "Profile & Customization", [
    "Premium Glassmorphism UI cards.",
    "Quick access to Support, Orders, and About.",
    "Secure Logout functionality.",
    "Editable profile details."
])

# Slide 16: Challenges & Solutions
add_bullet_slide(prs, "Challenges Faced", [
    "1. Firestore Indexing Errors -> Solved via Client-side sorting.",
    "2. Image Zoom Issues -> Fixed using BoxFit.contain.",
    "3. Missing Permissions -> Added Geolocation & Notification permissions.",
    "4. Build Errors -> Fixed Gradle desugaring dependencies."
])

# Slide 17: Future Roadmap
add_bullet_slide(prs, "Future Roadmap", [
    "Integration of UPI Payments (Razorpay/Paytm).",
    "AI-based Price Suggestions for Farmers.",
    "Delivery Partner App integration.",
    "Multi-language support (Hindi/Regional).",
    "Chat feature between Farmer and Buyer."
])

# Slide 18: Impact
add_bullet_slide(prs, "Project Impact", [
    "Empowers small-scale farmers.",
    "Provides healthy food options to ubran areas.",
    "Reduces food waste through direct selling.",
    "Scalable, modern, and user-friendly solution."
])

# Slide 19: Screen Gallery
add_image_slide(prs, "App Gallery", "[INSERT COLLAGE OF APP SCREENS: Profile, Splash, Login]")

# Slide 20: Conclusion
slide_end = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_end.shapes.title
subtitle = slide_end.placeholders[1]
title.text = "Thank You!"
subtitle.text = "Questions?\n\nContact: ritexstudios@farmigo.com"
title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

# --- Save Presentation ---
try:
    prs.save(OUTPUT_FILE)
    print(f"Successfully created presentation: {OUTPUT_FILE}")
except Exception as e:
    print(f"Error creating presentation: {e}")
